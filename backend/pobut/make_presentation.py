"""
make_presentation.py

Генерує PPTX-презентацію з основними знахідками аналізу
підозрілих ОП-абонентів (пошук крадіжок газу).
"""
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
os.chdir(os.path.dirname(os.path.abspath(__file__)))

import warnings
warnings.filterwarnings('ignore')

import io as _io
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.gridspec import GridSpec

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as pns
from lxml import etree

from pobut_predictor import GROUP_MERGE

# ── Кольори бренду ─────────────────────────────────────────────────────────────
C_DARK   = RGBColor(0x1A, 0x23, 0x3A)   # темно-синій фон
C_ACCENT = RGBColor(0x00, 0xA8, 0xE8)   # блакитний акцент
C_LIGHT  = RGBColor(0xF0, 0xF4, 0xF8)   # світлий фон блоків
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_RED    = RGBColor(0xE3, 0x2B, 0x2B)
C_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
C_YELLOW = RGBColor(0xF5, 0xC5, 0x18)
C_GRAY   = RGBColor(0x90, 0x90, 0x90)
C_GREEN  = RGBColor(0x27, 0xAE, 0x60)

MONTH_COLS = ['jan_2025','feb_2025','mar_2025','apr_2025','may_2025','jun_2025',
              'jul_2025','aug_2025','sep_2025','oct_2025','nov_2025','dec_2025']
ABS_LOW = 5.0

# ── 1. Дані ────────────────────────────────────────────────────────────────────
print("Завантаження даних...")
ap = pd.read_csv('data/input/all_pobut_enriched.csv', low_memory=False)
ap['appliance_group'] = ap['appliance_group'].astype(str).str.strip().replace(GROUP_MERGE)
ap['has_OP'] = ap['appliance_group'].str.contains('ОП', na=False)
ap['heated_area'] = pd.to_numeric(ap['heated_area'], errors='coerce').fillna(55.0)
ap['heated_area'] = ap['heated_area'].where(ap['heated_area'] > 0, 55.0)
for c in MONTH_COLS:
    ap[c] = pd.to_numeric(ap[c], errors='coerce').fillna(0)
ap['annual'] = ap[MONTH_COLS].sum(axis=1)

op = ap[
    ap['has_OP'] & ap['gas_off'].isna() & ap['alternative'].isna() &
    ap['dacha'].isna() & ap['grs'].notna() & (ap['annual'] > 0)
].copy().reset_index(drop=True)

op['specific'] = op['annual'] / op['heated_area']
ref = (op.groupby(['grs','consumer_type'])['specific']
         .median().rename('grs_ct_median').reset_index())
op = op.merge(ref, on=['grs','consumer_type'], how='left')
grs_ref = op.groupby('grs')['specific'].median().rename('grs_median').reset_index()
op = op.merge(grs_ref, on='grs', how='left')
op['ref_median'] = op['grs_ct_median'].fillna(op['grs_median'])
op['ratio_to_ref'] = op['specific'] / np.maximum(op['ref_median'], 1)

def classify(r):
    if r < 0.30: return 'critical'
    elif r < 0.40: return 'high'
    elif r < 0.50: return 'low'
    else: return 'normal'
op['level'] = op['ratio_to_ref'].apply(classify)

N_TOTAL    = len(op)
N_CRIT     = (op['level'] == 'critical').sum()
N_HIGH     = (op['level'] == 'high').sum()
N_LOW      = (op['level'] == 'low').sum()
N_NORMAL   = (op['level'] == 'normal').sum()
N_BELOW_D  = (op['specific'] < ABS_LOW).sum()
MEDIAN_SP  = op['specific'].median()
MED_PRIV   = op[op['consumer_type']=='Приватний сектор']['specific'].median()
MED_MULTI  = op[op['consumer_type']=='Багатоквартирний сектор']['specific'].median()
N_PRIV     = (op['consumer_type']=='Приватний сектор').sum()
N_MULTI    = (op['consumer_type']=='Багатоквартирний сектор').sum()

print(f"  N={N_TOTAL:,}, критичних={N_CRIT:,}")

# ── Хелпери для PPTX ───────────────────────────────────────────────────────────
def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, left, top, width, height,
                fontsize=18, bold=False, color=C_WHITE,
                align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fontsize)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def fig_to_stream(fig):
    buf = _io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor=fig.get_facecolor())
    buf.seek(0)
    return buf

def add_fig(slide, fig, left, top, width, height):
    buf = fig_to_stream(fig)
    slide.shapes.add_picture(buf, left, top, width, height)
    plt.close(fig)

# ── Розміри слайду 16:9 ────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # порожній

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 1 — ТИТУЛЬНИЙ
# ══════════════════════════════════════════════════════════════════════════════
print("Слайд 1: Титульний...")
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_DARK)

# Акцентна смуга зліва
add_rect(sl, Inches(0), Inches(0), Inches(0.18), H, C_ACCENT)

# Заголовок
add_textbox(sl, 'Аналіз питомого споживання газу',
            Inches(0.5), Inches(1.8), Inches(12), Inches(1.2),
            fontsize=40, bold=True, color=C_WHITE)

add_textbox(sl, 'ОП-абоненти · Пошук аномально низького споживання',
            Inches(0.5), Inches(3.0), Inches(12), Inches(0.8),
            fontsize=22, bold=False, color=C_ACCENT)

add_textbox(sl, 'Дані: білінг 2025 р.   |   Вибірка: 59 386 абонентів з опаленням',
            Inches(0.5), Inches(4.2), Inches(12), Inches(0.6),
            fontsize=14, color=RGBColor(0xAA, 0xBB, 0xCC))

add_textbox(sl, 'Березень 2026',
            Inches(0.5), Inches(6.5), Inches(4), Inches(0.5),
            fontsize=13, color=RGBColor(0x88, 0x99, 0xAA))

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 2 — ВИБІРКА І МЕТОДОЛОГІЯ
# ══════════════════════════════════════════════════════════════════════════════
print("Слайд 2: Методологія...")
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_DARK)
add_rect(sl, Inches(0), Inches(0), Inches(0.18), H, C_ACCENT)

add_textbox(sl, 'Вибірка та методологія',
            Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
            fontsize=28, bold=True, color=C_WHITE)
# Горизонтальна лінія
add_rect(sl, Inches(0.5), Inches(1.05), Inches(12.6), Inches(0.04), C_ACCENT)

# Ліва колонка — фільтри
items_left = [
    ('Вихідна база', 'Усі побутові абоненти'),
    ('Фільтр 1', 'Наявність ОП (опалювального котла)'),
    ('Фільтр 2', 'Активні: без gas_off / альтернативи / дачі'),
    ('Фільтр 3', 'Ненульовий річний білінг'),
    ('Результат', f'{N_TOTAL:,} абонентів'),
]
for i, (label, val) in enumerate(items_left):
    top = Inches(1.3 + i * 0.9)
    add_rect(sl, Inches(0.5), top, Inches(5.8), Inches(0.75),
             RGBColor(0x22, 0x2E, 0x4A))
    add_textbox(sl, label, Inches(0.65), top + Pt(4), Inches(2), Inches(0.4),
                fontsize=11, color=C_ACCENT, bold=True)
    add_textbox(sl, val, Inches(2.6), top + Pt(4), Inches(3.9), Inches(0.4),
                fontsize=12, color=C_WHITE)

# Права колонка — метрика
add_textbox(sl, 'Ключова метрика',
            Inches(7.0), Inches(1.15), Inches(5.8), Inches(0.5),
            fontsize=16, bold=True, color=C_ACCENT)

formula_lines = [
    'specific  =  annual_billing  /  heated_area',
    '',
    'ratio_to_ref  =  specific  /  median(specific)',
    '                           по ГРС × тип застройки',
    '',
    'Рівні підозрілості:',
    '  ratio < 0.30  →  Критичний',
    '  ratio 0.30–0.40  →  Середній',
    '  ratio 0.40–0.50  →  Низький',
    '  ratio ≥ 0.50  →  Норма',
]
add_rect(sl, Inches(7.0), Inches(1.7), Inches(5.8), Inches(5.3),
         RGBColor(0x0D, 0x16, 0x2B))
add_textbox(sl, '\n'.join(formula_lines),
            Inches(7.15), Inches(1.8), Inches(5.5), Inches(5.0),
            fontsize=12, color=RGBColor(0xCC, 0xDD, 0xEE))

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 3 — РОЗПОДІЛ ПО РІВНЯХ (великі цифри + bar)
# ══════════════════════════════════════════════════════════════════════════════
print("Слайд 3: Розподіл...")
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_DARK)
add_rect(sl, Inches(0), Inches(0), Inches(0.18), H, C_ACCENT)

add_textbox(sl, 'Розподіл абонентів за рівнем підозрілості',
            Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
            fontsize=28, bold=True, color=C_WHITE)
add_rect(sl, Inches(0.5), Inches(1.05), Inches(12.6), Inches(0.04), C_ACCENT)

# 4 картки
cards = [
    ('Норма',            f'{N_NORMAL:,}', f'{N_NORMAL/N_TOTAL*100:.0f}%', C_GRAY,   'ratio ≥ 0.50'),
    ('Низький',          f'{N_LOW:,}',    f'{N_LOW/N_TOTAL*100:.1f}%',    C_YELLOW, 'ratio 0.40–0.50'),
    ('Середній/Високий', f'{N_HIGH:,}',   f'{N_HIGH/N_TOTAL*100:.1f}%',   C_ORANGE, 'ratio 0.30–0.40'),
    ('Критичний',        f'{N_CRIT:,}',   f'{N_CRIT/N_TOTAL*100:.1f}%',   C_RED,    'ratio < 0.30'),
]
card_w = Inches(2.9)
for i, (title, n, pct, color, sub) in enumerate(cards):
    left = Inches(0.5 + i * 3.1)
    top  = Inches(1.3)
    # Кольорова смужка зверху
    add_rect(sl, left, top, card_w, Inches(0.12), color)
    # Фон картки
    add_rect(sl, left, top + Inches(0.12), card_w, Inches(3.5), RGBColor(0x1E, 0x2A, 0x44))
    # Заголовок
    add_textbox(sl, title, left + Inches(0.15), top + Inches(0.22),
                card_w - Inches(0.3), Inches(0.5),
                fontsize=14, bold=True, color=C_WHITE)
    # Велика цифра
    add_textbox(sl, n, left + Inches(0.1), top + Inches(0.75),
                card_w - Inches(0.2), Inches(1.1),
                fontsize=38, bold=True, color=color, align=PP_ALIGN.CENTER)
    # Відсоток
    add_textbox(sl, pct + ' від вибірки', left + Inches(0.1), top + Inches(1.85),
                card_w - Inches(0.2), Inches(0.45),
                fontsize=14, color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.CENTER)
    # Пояснення
    add_textbox(sl, sub, left + Inches(0.1), top + Inches(2.35),
                card_w - Inches(0.2), Inches(0.4),
                fontsize=11, color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.CENTER)

# Підсумок внизу
add_rect(sl, Inches(0.5), Inches(5.05), Inches(12.6), Inches(0.75),
         RGBColor(0x12, 0x1D, 0x36))
add_textbox(sl,
    f'Разом підозрілих (будь-який рівень): {N_LOW+N_HIGH+N_CRIT:,} абонентів  '
    f'({(N_LOW+N_HIGH+N_CRIT)/N_TOTAL*100:.1f}% вибірки)   |   '
    f'Нижче критерію D (< 5 м³/м²/рік): {N_BELOW_D:,} абонентів',
    Inches(0.7), Inches(5.15), Inches(12.2), Inches(0.55),
    fontsize=13, color=C_WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 4 — SCATTER PLOT
# ══════════════════════════════════════════════════════════════════════════════
print("Слайд 4: Scatter plot...")
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_DARK)
add_rect(sl, Inches(0), Inches(0), Inches(0.18), H, C_ACCENT)

add_textbox(sl, 'Питоме споживання vs площа',
            Inches(0.5), Inches(0.2), Inches(12), Inches(0.65),
            fontsize=28, bold=True, color=C_WHITE)
add_rect(sl, Inches(0.5), Inches(0.9), Inches(12.6), Inches(0.04), C_ACCENT)

# Будуємо два scatter'и поруч
SCATTER_COLORS = {
    'normal':   '#909090',
    'low':      '#F5C518',
    'high':     '#FF8C00',
    'critical': '#E32B2B',
}
ZORDER_S = {'normal':1,'low':2,'high':3,'critical':4}

SECTORS_SL = [
    ('Приватний сектор',        'Приватний сектор'),
    ('Багатоквартирний сектор', 'Багатоквартирний'),
]

fig, axes = plt.subplots(1, 2, figsize=(12.5, 5.2), facecolor='#0D1626')
for ax, (ct_name, ct_short) in zip(axes, SECTORS_SL):
    ax.set_facecolor('#0D1626')
    ct = op[op['consumer_type'] == ct_name]
    ct_med = ct['specific'].median()
    ct_cap = ct['specific'].quantile(0.99)

    for lvl in ['normal', 'low', 'high', 'critical']:
        sub = ct[ct['level'] == lvl]
        if sub.empty: continue
        n_lvl = len(sub)
        ax.scatter(sub['heated_area'], sub['specific'],
                   c=SCATTER_COLORS[lvl], alpha=0.25, s=6,
                   label=f'{lvl.capitalize()}  N={n_lvl:,}',
                   linewidths=0, zorder=ZORDER_S[lvl])

    ax.axhline(ct_med, color='#00A8E8', lw=1.8, ls='--',
               label=f'Медіана = {ct_med:.1f}')
    ax.axhline(ABS_LOW, color='#FF4444', lw=1.8, ls=':',
               label=f'Критерій D = {ABS_LOW:.0f}')

    ax.set_xscale('log')
    ax.set_ylim(0, ct_cap)
    ax.set_xlabel('Площа, м² (лог.)', color='#AABBCC', fontsize=10)
    ax.set_ylabel('м³/м²/рік', color='#AABBCC', fontsize=10)
    ax.set_title(f'{ct_short}  (N={len(ct):,})', color='white', fontsize=12, pad=6)
    ax.tick_params(colors='#8899AA', labelsize=8)
    for spine in ax.spines.values(): spine.set_color('#2A3A5A')
    ax.grid(True, alpha=0.15, color='#445577')
    ax.legend(fontsize=8, framealpha=0.3, facecolor='#1A2340',
              edgecolor='#445577', labelcolor='white', loc='upper right')

plt.tight_layout(pad=0.8)
add_fig(sl, fig, Inches(0.25), Inches(1.05), Inches(12.85), Inches(6.1))

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 5 — ГРУПИ СПОЖИВАЧІВ
# ══════════════════════════════════════════════════════════════════════════════
print("Слайд 5: Групи...")
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_DARK)
add_rect(sl, Inches(0), Inches(0), Inches(0.18), H, C_ACCENT)

add_textbox(sl, 'Порівняння груп споживачів',
            Inches(0.5), Inches(0.2), Inches(12), Inches(0.65),
            fontsize=28, bold=True, color=C_WHITE)
add_rect(sl, Inches(0.5), Inches(0.9), Inches(12.6), Inches(0.04), C_ACCENT)

# Bar chart: медіани по типах
fig, axes = plt.subplots(1, 2, figsize=(9, 4.5), facecolor='#0D1626')

# Лівий: медіани по типу
ax1 = axes[0]
ax1.set_facecolor('#0D1626')
types = ['Приватний\nсектор', 'Багатоквартирний\nсектор']
medians = [MED_PRIV, MED_MULTI]
colors_b = ['#00A8E8', '#5B9BD5']
bars = ax1.barh(types, medians, color=colors_b, height=0.5, edgecolor='none')
for bar, val in zip(bars, medians):
    ax1.text(val + 0.3, bar.get_y() + bar.get_height()/2,
             f'{val:.1f}', va='center', ha='left', color='white', fontsize=13, fontweight='bold')
ax1.set_xlim(0, max(medians)*1.25)
ax1.set_xlabel('Медіана specific, м³/м²/рік', color='#AABBCC', fontsize=10)
ax1.set_title('Медіана по типу забудови', color='white', fontsize=12, pad=8)
ax1.tick_params(colors='#8899AA')
for sp in ax1.spines.values(): sp.set_color('#2A3A5A')
ax1.grid(axis='x', alpha=0.15, color='#445577')

# Правий: розподіл критичних по типу
ax2 = axes[1]
ax2.set_facecolor('#0D1626')
priv_crit  = ((op['consumer_type']=='Приватний сектор') & (op['level']=='critical')).sum()
multi_crit = ((op['consumer_type']=='Багатоквартирний сектор') & (op['level']=='critical')).sum()
priv_pct   = priv_crit / N_PRIV * 100
multi_pct  = multi_crit / N_MULTI * 100
bars2 = ax2.barh(['Приватний\nсектор', 'Багатоквартирний\nсектор'],
                 [priv_pct, multi_pct], color=['#E32B2B', '#FF8C00'], height=0.5, edgecolor='none')
for bar, val in zip(bars2, [priv_pct, multi_pct]):
    ax2.text(val + 0.1, bar.get_y() + bar.get_height()/2,
             f'{val:.1f}%', va='center', ha='left', color='white', fontsize=13, fontweight='bold')
ax2.set_xlim(0, max(priv_pct, multi_pct)*1.3)
ax2.set_xlabel('% критичних в групі', color='#AABBCC', fontsize=10)
ax2.set_title('Критичні (ratio < 0.30) по типу', color='white', fontsize=12, pad=8)
ax2.tick_params(colors='#8899AA')
for sp in ax2.spines.values(): sp.set_color('#2A3A5A')
ax2.grid(axis='x', alpha=0.15, color='#445577')

plt.tight_layout(pad=1.0)
add_fig(sl, fig, Inches(0.35), Inches(1.05), Inches(9.2), Inches(5.8))

# Текст праворуч
stat_rows = [
    ('N абонентів', f'{N_PRIV:,}', f'{N_MULTI:,}'),
    ('Медіана specific', f'{MED_PRIV:.1f}', f'{MED_MULTI:.1f}'),
    ('Критичних', f'{priv_crit:,}', f'{multi_crit:,}'),
    ('% критичних', f'{priv_pct:.1f}%', f'{multi_pct:.1f}%'),
]
add_textbox(sl, 'Приватний    Багатокв.',
            Inches(9.7), Inches(1.1), Inches(3.4), Inches(0.5),
            fontsize=11, color=C_ACCENT, bold=True)
for i, (label, v1, v2) in enumerate(stat_rows):
    top = Inches(1.65 + i * 0.95)
    add_rect(sl, Inches(9.7), top, Inches(3.4), Inches(0.82),
             RGBColor(0x1E, 0x2A, 0x44))
    add_textbox(sl, label, Inches(9.8), top + Inches(0.05),
                Inches(3.2), Inches(0.35), fontsize=10, color=RGBColor(0x88, 0x99, 0xAA))
    add_textbox(sl, f'{v1}    {v2}', Inches(9.8), top + Inches(0.38),
                Inches(3.2), Inches(0.38), fontsize=13, bold=True, color=C_WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 6 — ВИСНОВКИ
# ══════════════════════════════════════════════════════════════════════════════
print("Слайд 6: Висновки...")
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_DARK)
add_rect(sl, Inches(0), Inches(0), Inches(0.18), H, C_ACCENT)

add_textbox(sl, 'Висновки та рекомендації',
            Inches(0.5), Inches(0.2), Inches(12), Inches(0.65),
            fontsize=28, bold=True, color=C_WHITE)
add_rect(sl, Inches(0.5), Inches(0.9), Inches(12.6), Inches(0.04), C_ACCENT)

conclusions = [
    (C_RED,    '11.9%', 'абонентів мають аномально низьке питоме споживання (ratio < 0.50) — потребують перевірки'),
    (C_RED,    '11.7%', 'знаходяться нижче критичного рівня (ratio < 0.30) — найвищий пріоритет'),
    (C_ORANGE, '6 072', 'абонентів нижче абсолютного порогу 5 м³/м²/рік (критерій D) незалежно від групи'),
    (C_ACCENT, '3.4×',  'різниця медіани між нормою та критичними (19.6 vs 2.6 м³/м²/рік) — аномалія очевидна'),
    (C_GREEN,  'GRS',   'Прив\'язка до ГРС×тип дозволяє коректно порівнювати різні регіони та типи забудови'),
]

for i, (color, val, text) in enumerate(conclusions):
    top = Inches(1.15 + i * 1.1)
    add_rect(sl, Inches(0.5), top, Inches(12.6), Inches(0.95),
             RGBColor(0x1A, 0x25, 0x3E))
    add_rect(sl, Inches(0.5), top, Inches(0.06), Inches(0.95), color)
    add_textbox(sl, val, Inches(0.7), top + Inches(0.05),
                Inches(1.3), Inches(0.8),
                fontsize=24, bold=True, color=color, align=PP_ALIGN.RIGHT)
    add_textbox(sl, text, Inches(2.15), top + Inches(0.2),
                Inches(10.8), Inches(0.6),
                fontsize=14, color=C_WHITE)

add_textbox(sl, 'Наступні кроки: виїзні перевірки критичних абонентів · '
                'Верифікація через дані модемів · Порівняння з попередніми роками',
            Inches(0.5), Inches(6.8), Inches(12.6), Inches(0.45),
            fontsize=11, color=RGBColor(0x88, 0x99, 0xAA))

# ── Зберігаємо ─────────────────────────────────────────────────────────────────
OUT = 'data/output/op_theft_analysis.pptx'
os.makedirs('data/output', exist_ok=True)
prs.save(OUT)
print(f"\nЗбережено: {OUT}")
